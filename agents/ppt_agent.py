from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
import matplotlib.pyplot as plt
import datetime


def generate_pie_chart(labels, values, chart_path="chart.png"):
    fig, ax = plt.subplots()
    wedges, texts = ax.pie(
        values, labels=None, autopct=None, startangle=90, wedgeprops=dict(width=0.4)
    )
    ax.axis('equal')
    plt.legend(wedges, [f"{l}: {v}%" for l, v in zip(labels, values)], loc="lower center", ncol=2, bbox_to_anchor=(0.5, -0.2))
    plt.savefig(chart_path, bbox_inches='tight')
    plt.close()

def update_template(summary_data, template_path="template.pptx", output_path="monthly_status_updated.pptx"):
    prs = Presentation(template_path)
    slide = prs.slides[0]  # You can change this if multiple slides
    # Set slide title
    month_year = datetime.datetime.now().strftime("%B %Y")
    title_text = f"Monthly Status Update - {month_year}"
    # Try to find a title shape 
    title_set = False
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX and shape.has_text_frame:
            if "{{Title}}" in shape.text_frame.text:
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                p.text = title_text
                p.font.size = Pt(32)
                title_set = True
                break
    # Replace placeholders
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX and shape.has_text_frame:
            text = shape.text_frame.text
            for section in ["Current Status", "Key Themes", "Risks", "Blockers"]:
                placeholder = f"{{{{{section}}}}}"
                if placeholder in text:
                    shape.text_frame.clear()
                    points = summary_data["bullets"].get(section, [])
                    # If points is a string, convert to a single-item list
                    if isinstance(points, str):
                        points = [points]
                    for point in points:
                        p = shape.text_frame.add_paragraph()
                        p.text = f"{point}"
                        p.font.size = Pt(14)

    # Generate and insert pie chart
    if "chart" in summary_data and "labels" in summary_data["chart"] and "values" in summary_data["chart"]:
        chart_path = "chart.png"
        labels = summary_data["chart"]["labels"]
        values = summary_data["chart"]["values"]

        generate_pie_chart(labels, values, chart_path)

        # Insert image at a fixed position
        slide.shapes.add_picture(chart_path, Inches(5), Inches(2.5), width=Inches(4), height=Inches(3))
    else:
        print("Warning: No chart data found in summary_data. Skipping chart generation.")

    prs.save(output_path)
    return {
        "success": True,
        "message": f"PowerPoint presentation updated and saved to {output_path}"
    }